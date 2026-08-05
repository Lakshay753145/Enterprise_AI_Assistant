"""Document parsing.

Deliberately lightweight: pypdf for PDFs, python-docx for Word, openpyxl for
Excel, plain reads for text formats. No ML models, no multi-gigabyte downloads,
no OCR pipeline - parsing starts instantly and a 50-page PDF takes about a
second.

The trade-off, stated plainly: complex tables and scanned (image-only) pages
extract poorly or not at all. For born-digital policy and procedure documents -
which is the bulk of an internal knowledge base - the text comes out fine.

Text is captured **per page**, so every chunk keeps a real page number and
citations point somewhere a person can actually turn to.
"""

from __future__ import annotations

import hashlib
import re
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from backend.core.exceptions import IngestionError, UnsupportedFileTypeError
from backend.core.logging_config import ingestion_logger


@dataclass
class PageText:
    """One page (or sheet, or slide) of extracted text."""

    page_number: int
    text: str


@dataclass
class ParsedDocument:
    """Normalised parse output, ready for chunking."""

    pages: list[PageText]
    title: str | None
    page_count: int
    file_hash: str
    size_bytes: int
    parse_seconds: float
    metadata: dict[str, Any] = field(default_factory=dict)

    @property
    def text(self) -> str:
        return "\n\n".join(p.text for p in self.pages)

    @property
    def char_count(self) -> int:
        return sum(len(p.text) for p in self.pages)


def compute_file_hash(path: Path) -> str:
    """SHA-256 of the file bytes, streamed so large files do not blow memory."""
    digest = hashlib.sha256()
    with Path(path).open("rb") as fh:
        for block in iter(lambda: fh.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def warmup_converter() -> None:
    """No-op. Kept so callers do not need to care which parser is in use."""
    return None


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def parse_document(path: Path | str) -> ParsedDocument:
    """Extract text from a file.

    Raises :class:`IngestionError` when the result is unusable. An empty parse
    is treated as a failure rather than an empty document, because silently
    ingesting zero chunks looks like success and then nothing is answerable.
    """
    path = Path(path)
    if not path.is_file():
        raise IngestionError(f"File not found: {path}")

    suffix = path.suffix.lower()
    handler = _HANDLERS.get(suffix)
    if handler is None:
        raise UnsupportedFileTypeError(
            f"{suffix or 'This file type'} is not supported. Supported: "
            + ", ".join(sorted(_HANDLERS))
        )

    size_bytes = path.stat().st_size
    file_hash = compute_file_hash(path)
    started = time.perf_counter()

    ingestion_logger.info(f"Parsing {path.name} ({size_bytes / 1_048_576:.1f} MB)...")

    try:
        pages, extra = handler(path)
    except IngestionError:
        raise
    except Exception as exc:
        ingestion_logger.exception(f"Parse failed for {path.name}")
        raise IngestionError(f"Could not read {path.name}: {exc}") from exc

    pages = [PageText(p.page_number, _tidy(p.text)) for p in pages]
    pages = [p for p in pages if p.text.strip()]

    if not pages:
        raise IngestionError(
            f"{path.name} contained no extractable text. If this is a scanned "
            f"PDF, the pages are images - run it through OCR first, or supply "
            f"a text-based version."
        )

    elapsed = time.perf_counter() - started
    title = _extract_title(pages, path)

    ingestion_logger.info(
        f"Parsed {path.name}: {len(pages)} pages, "
        f"{sum(len(p.text) for p in pages):,} chars, {elapsed:.2f}s"
    )

    return ParsedDocument(
        pages=pages,
        title=title,
        page_count=len(pages),
        file_hash=file_hash,
        size_bytes=size_bytes,
        parse_seconds=round(elapsed, 3),
        metadata={"source_filename": path.name, "parser": suffix.lstrip("."), **extra},
    )


# ---------------------------------------------------------------------------
# Per-format handlers
# ---------------------------------------------------------------------------

def _parse_pdf(path: Path) -> tuple[list[PageText], dict[str, Any]]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))

    if reader.is_encrypted:
        # Many "protected" PDFs use an empty owner password and open fine.
        try:
            reader.decrypt("")
        except Exception as exc:
            raise IngestionError(
                f"{path.name} is password-protected. Remove the password and "
                f"re-upload."
            ) from exc

    pages: list[PageText] = []
    for index, page in enumerate(reader.pages, start=1):
        try:
            pages.append(PageText(index, page.extract_text() or ""))
        except Exception as exc:
            ingestion_logger.warning(f"{path.name} page {index} unreadable: {exc}")

    info = reader.metadata or {}
    extra: dict[str, Any] = {}
    if getattr(info, "title", None):
        extra["pdf_title"] = str(info.title)[:300]

    empty = sum(1 for p in pages if not p.text.strip())
    if pages and empty == len(pages):
        raise IngestionError(
            f"{path.name} has {len(pages)} pages but no selectable text - it is "
            f"almost certainly a scan. Run OCR on it first."
        )
    if empty:
        ingestion_logger.warning(
            f"{path.name}: {empty} of {len(pages)} pages had no text (likely images)"
        )

    return pages, extra


def _parse_docx(path: Path) -> tuple[list[PageText], dict[str, Any]]:
    from docx import Document as DocxDocument

    document = DocxDocument(str(path))
    parts: list[str] = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        # Preserve heading level as markdown so the chunker can see structure.
        style = (paragraph.style.name or "").lower() if paragraph.style else ""
        if style.startswith("heading"):
            level = "".join(c for c in style if c.isdigit()) or "2"
            parts.append(f"{'#' * min(int(level), 6)} {text}")
        else:
            parts.append(text)

    for table in document.tables:
        rows = [
            "| " + " | ".join(cell.text.strip() for cell in row.cells) + " |"
            for row in table.rows
        ]
        if len(rows) >= 2:
            separator = "|" + "---|" * len(table.rows[0].cells)
            rows.insert(1, separator)
        parts.append("\n".join(rows))

    # Word has no fixed pagination we can see; treat it as one logical page.
    return [PageText(1, "\n\n".join(parts))], {}


def _parse_xlsx(path: Path) -> tuple[list[PageText], dict[str, Any]]:
    from openpyxl import load_workbook

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    pages: list[PageText] = []

    for index, sheet in enumerate(workbook.worksheets, start=1):
        rows: list[str] = [f"# {sheet.title}"]
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append("| " + " | ".join(cells) + " |")
        if len(rows) > 1:
            pages.append(PageText(index, "\n".join(rows)))

    workbook.close()
    return pages, {}


def _parse_pptx(path: Path) -> tuple[list[PageText], dict[str, Any]]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    pages: list[PageText] = []

    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
        if parts:
            pages.append(PageText(index, "\n\n".join(parts)))

    return pages, {}


def _parse_text(path: Path) -> tuple[list[PageText], dict[str, Any]]:
    raw = _read_text(path)
    return [PageText(1, raw)], {}


def _parse_html(path: Path) -> tuple[list[PageText], dict[str, Any]]:
    raw = _read_text(path)
    raw = re.sub(r"(?is)<(script|style).*?</\1>", " ", raw)
    raw = re.sub(r"(?i)<br\s*/?>|</p>|</div>|</tr>", "\n", raw)
    raw = re.sub(r"(?i)<h([1-6])[^>]*>", lambda m: "\n" + "#" * int(m.group(1)) + " ", raw)
    raw = re.sub(r"<[^>]+>", " ", raw)
    raw = (
        raw.replace("&nbsp;", " ")
        .replace("&amp;", "&")
        .replace("&lt;", "<")
        .replace("&gt;", ">")
        .replace("&quot;", '"')
    )
    return [PageText(1, raw)], {}


def _read_text(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_bytes().decode("utf-8", errors="replace")


_HANDLERS = {
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".xlsx": _parse_xlsx,
    ".pptx": _parse_pptx,
    ".txt": _parse_text,
    ".md": _parse_text,
    ".html": _parse_html,
    ".htm": _parse_html,
}


# ---------------------------------------------------------------------------
# Clean-up
# ---------------------------------------------------------------------------

def _tidy(text: str) -> str:
    """Repair the usual PDF extraction artefacts."""
    if not text:
        return ""

    # Words split across a line break by hyphenation: "solu-\ntion" -> "solution"
    text = re.sub(r"(\w)-\s*\n\s*(\w)", r"\1\2", text)
    # Collapse runs of spaces, but keep newlines - they carry list structure.
    text = re.sub(r"[ \t]{2,}", " ", text)
    # Three or more blank lines become a paragraph break.
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Strip trailing whitespace per line.
    text = "\n".join(line.rstrip() for line in text.splitlines())
    return text.strip()


def _extract_title(pages: list[PageText], path: Path) -> str | None:
    """First markdown H1, else the first substantial line, else the filename."""
    if pages:
        for line in pages[0].text.splitlines():
            stripped = line.strip()
            if stripped.startswith("# "):
                return stripped[2:].strip()[:300]
        for line in pages[0].text.splitlines():
            stripped = line.strip()
            if 10 <= len(stripped) <= 120 and not stripped.endswith("."):
                return stripped[:300]

    return path.stem.replace("_", " ").replace("-", " ").strip()[:300] or None
